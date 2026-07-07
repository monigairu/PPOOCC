"""
補足資料（Excel/PPTX）から画像を抽出し Gemini でキャプションを生成するスクリプト

実行方法:
    uv run python scripts/preliminary_review/generate_supplement_captions.py           # 全ファイル処理
    uv run python scripts/preliminary_review/generate_supplement_captions.py --file 東北電力_補足.xlsx
    uv run python scripts/preliminary_review/generate_supplement_captions.py --dry-run # Gemini呼び出しなし（画像抽出のみ確認）

出力:
    data/knowledge/supplement_captions/{source_file_stem}.json
    （Vertex AI Search 投入前に内容を確認・修正できる中間ファイル）

出力JSONの1件の構造:
    {
        "id":               "東北電力_東北電力_補足_001",
        "caption":          "PPパネルが完全に撤去されており工事完了状態",
        "utility_name":     "東北電力",
        "fee_type":         "解体費",
        "source_file":      "東北電力_補足.xlsx",
        "construction_name":"○○建屋解体工事",
        "context_text":     "撤去後",
        "original_format":  "excel"
    }
"""
from __future__ import annotations

import argparse
import json
import logging
import re
import sys
from pathlib import Path
from typing import Any

sys.path.insert(0, str(Path(__file__).parent.parent.parent))

from google import genai
from google.genai import types

logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")
logger = logging.getLogger(__name__)

_SUPPLEMENT_DIR = Path("data/knowledge/supplement")
_CAPTION_DIR    = Path("data/knowledge/supplement_captions")
_MODEL          = "gemini-2.5-flash"

# ファイル名から電力会社名を推定するパターン（例：「東北電力_補足.xlsx」→「東北電力」）
_UTILITY_RE = re.compile(r"^([^_]+)")


def _get_genai_client() -> genai.Client:
    return genai.Client(vertexai=True)


# ── キャプション生成 ──────────────────────────────────────────────────────────

def _build_prompt(context_text: str, construction_name: str, original_format: str, work_area: str = "") -> str:
    if original_format == "pptx":
        area = f"（{work_area}エリア）" if work_area else ""
        return (
            f"建屋平面図の工事進捗図です。「{construction_name}」{area}を示しています。"
            "図中の色分けや工事の進捗状態を日本語で具体的に説明してください。"
            "説明は2〜4文程度で、工事の完了状況・作業エリア・特記事項を含めてください。"
        )
    label = f"「{context_text}」と書かれた枠の" if context_text else ""
    return (
        f"「{construction_name}」の補足資料です。{label}写真です。"
        "工事の状態や作業内容を日本語で具体的に説明してください。"
        "説明は2〜4文程度で、工事の完了状況・撤去物・作業の進捗を含めてください。"
    )


def _generate_caption(
    client: genai.Client,
    img_bytes: bytes,
    mime_type: str,
    context_text: str,
    construction_name: str,
    original_format: str,
    work_area: str = "",
    dry_run: bool = False,
) -> str:
    if dry_run:
        return f"[DRY-RUN] {context_text or construction_name} のキャプション（未生成）"

    prompt = _build_prompt(context_text, construction_name, original_format, work_area)
    try:
        response = client.models.generate_content(
            model=_MODEL,
            contents=[
                types.Part.from_bytes(data=img_bytes, mime_type=mime_type),
                types.Part.from_text(text=prompt),
            ],
            config=types.GenerateContentConfig(temperature=0.0),
        )
        return response.text.strip()
    except Exception as e:
        logger.warning("キャプション生成エラー（スキップ）: %s", e)
        return ""


# ── Excel 処理 ────────────────────────────────────────────────────────────────

def _extract_from_excel(file_path: Path) -> list[dict[str, Any]]:
    """
    Excel から画像バイト列と周辺テキストを抽出する。

    想定フォーマット:
        A1: 工事名
        各シート: 「撤去前」「撤去後」等のラベルセルと画像が混在
    """
    import openpyxl

    records = []
    wb = openpyxl.load_workbook(file_path)

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        construction_name = ""
        if ws.cell(1, 1).value:
            construction_name = str(ws.cell(1, 1).value).strip()

        # シート内の全セルテキストを row/col インデックスで収集（画像のアンカー参照用）
        cell_texts: dict[tuple[int, int], str] = {}
        for row in ws.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str) and cell.value.strip():
                    cell_texts[(cell.row, cell.column)] = cell.value.strip()

        if not hasattr(ws, "_images") or not ws._images:
            continue

        for idx, img_obj in enumerate(ws._images, 1):
            try:
                img_bytes = img_obj._data()
            except Exception:
                continue

            # 画像のアンカー位置（左上セル）を取得
            anchor = img_obj.anchor
            try:
                row = anchor._from.row + 1   # 0-indexed → 1-indexed
                col = anchor._from.col + 1
            except AttributeError:
                row, col = 1, 1

            # 画像の左・上・右のセルからコンテキストテキストを探す
            context_text = ""
            for dr, dc in [(0, -1), (-1, 0), (0, 1), (1, 0)]:
                t = cell_texts.get((row + dr, col + dc), "")
                if t:
                    context_text = t
                    break

            records.append({
                "_img_bytes": img_bytes,
                "_mime_type": "image/png",
                "source_file": file_path.name,
                "sheet_name": sheet_name,
                "construction_name": construction_name or sheet_name,
                "context_text": context_text,
                "original_format": "excel",
                "image_index": idx,
            })

    return records


# ── PPTX 処理 ─────────────────────────────────────────────────────────────────

def _extract_from_pptx(file_path: Path) -> list[dict[str, Any]]:
    """
    PPTX から画像バイト列とスライドテキストを抽出する。

    想定フォーマット:
        各スライド: 建屋平面図 + 工事名・作業エリア名テキスト
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    records = []
    prs = Presentation(file_path)

    for slide_idx, slide in enumerate(prs.slides, 1):
        texts = []
        images = []

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append({
                    "bytes": shape.image.blob,
                    "mime_type": f"image/{shape.image.ext.lower()}",
                })
            elif shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)

        if not images:
            continue

        construction_name = texts[0] if texts else f"スライド{slide_idx}"
        work_area = texts[1] if len(texts) > 1 else ""

        for img_idx, img in enumerate(images, 1):
            records.append({
                "_img_bytes": img["bytes"],
                "_mime_type": img["mime_type"],
                "source_file": file_path.name,
                "sheet_name": f"slide_{slide_idx:03d}",
                "construction_name": construction_name,
                "context_text": work_area,
                "original_format": "pptx",
                "image_index": img_idx,
            })

    return records


# ── ファイル処理のエントリ ────────────────────────────────────────────────────

def process_file(file_path: Path, client: genai.Client, dry_run: bool) -> list[dict]:
    suffix = file_path.suffix.lower()
    if suffix == ".xlsx":
        raw_records = _extract_from_excel(file_path)
    elif suffix == ".pptx":
        raw_records = _extract_from_pptx(file_path)
    else:
        logger.warning("未対応のファイル形式をスキップ: %s", file_path.name)
        return []

    if not raw_records:
        logger.info("  画像なし: %s", file_path.name)
        return []

    # ファイル名の先頭部分から電力会社名を推定
    m = _UTILITY_RE.match(file_path.stem)
    utility_name = m.group(1) if m else file_path.stem

    # ファイル名のベース部分（拡張子なし・特殊文字除去）をID用に正規化
    safe_stem = re.sub(r"[^\w]", "_", file_path.stem)[:40]

    results = []
    for rec in raw_records:
        img_bytes = rec.pop("_img_bytes")
        mime_type = rec.pop("_mime_type")

        caption = _generate_caption(
            client=client,
            img_bytes=img_bytes,
            mime_type=mime_type,
            context_text=rec["context_text"],
            construction_name=rec["construction_name"],
            original_format=rec["original_format"],
            work_area=rec.get("context_text", ""),
            dry_run=dry_run,
        )
        if not caption:
            continue

        doc_id = f"{utility_name}_{safe_stem}_{rec['image_index']:03d}"
        doc_id = re.sub(r"[^\w-]", "_", doc_id)[:128]

        results.append({
            "id":                doc_id,
            "caption":           caption,
            "utility_name":      utility_name,
            "fee_type":          "",          # ファイル名・シート名から推定できる場合は上書き
            "source_file":       rec["source_file"],
            "sheet_name":        rec["sheet_name"],
            "construction_name": rec["construction_name"],
            "context_text":      rec["context_text"],
            "original_format":   rec["original_format"],
        })
        logger.info("  [%s] %s → %d文字のキャプション生成",
                    doc_id, rec["context_text"] or rec["construction_name"], len(caption))

    return results


def main() -> None:
    parser = argparse.ArgumentParser(description="補足資料キャプション生成スクリプト")
    parser.add_argument("--file", help="処理対象ファイル名（省略時は全ファイル）")
    parser.add_argument("--dry-run", action="store_true", help="Gemini呼び出しなし（画像抽出の確認のみ）")
    args = parser.parse_args()

    if not _SUPPLEMENT_DIR.exists():
        print(f"エラー: 補足資料ディレクトリが見つかりません: {_SUPPLEMENT_DIR}")
        sys.exit(1)

    if args.file:
        target_files = [_SUPPLEMENT_DIR / args.file]
        missing = [f for f in target_files if not f.exists()]
        if missing:
            print(f"エラー: ファイルが見つかりません: {missing[0]}")
            sys.exit(1)
    else:
        target_files = sorted(
            list(_SUPPLEMENT_DIR.glob("*.xlsx")) + list(_SUPPLEMENT_DIR.glob("*.pptx"))
        )
        if not target_files:
            print(f"補足資料ファイルが見つかりません: {_SUPPLEMENT_DIR}")
            print("data/knowledge/supplement/ に .xlsx または .pptx ファイルを配置してから実行してください")
            sys.exit(0)

    _CAPTION_DIR.mkdir(parents=True, exist_ok=True)
    client = _get_genai_client()

    total = 0
    for file_path in target_files:
        logger.info("処理中: %s", file_path.name)
        records = process_file(file_path, client, dry_run=args.dry_run)

        out_path = _CAPTION_DIR / f"{file_path.stem}.json"
        with open(out_path, "w", encoding="utf-8") as f:
            json.dump(records, f, ensure_ascii=False, indent=2)

        logger.info("  → %d件 を %s に保存", len(records), out_path)
        total += len(records)

    print(f"\n完了: 合計 {total} 件のキャプションを生成しました")
    if not args.dry_run:
        print("次のステップ: uv run python scripts/preliminary_review/ingest_knowledge.py --target supplement")


if __name__ == "__main__":
    main()
